"""多格式附件文本抽取。

设计原则（与项目现有 OCR / vision 模块保持一致）：
- 输入：原始字节 + 文件名（用扩展名分发解析器）。
- 输出：尽量短的纯文本，控制单文件抽取上限以保护 prompt token 预算。
- 抽取失败不抛异常，返回带 ``[抽取失败: ...]`` 前缀的占位文本。
- 解析依赖延迟导入：未安装的依赖只在真的遇到对应文件时才报错，
  并把错误以可读的中文形式返回，而不是把整个请求拖死。

接入方：``server.py`` 的 ``/api/chat/stream``，把抽取文本拼到用户消息前。
原文件本身不持久化，与图片 OCR 流程一致。
"""
from __future__ import annotations

import io
import os
import zipfile
from typing import Tuple

# 单个附件抽取出的纯文本上限（字符）。超过会截断并提示。
PER_FILE_CHAR_LIMIT = 8000
# ZIP 内最多处理多少个文件，超出忽略。
ZIP_MAX_INNER_FILES = 50
# 嵌套 ZIP 递归深度上限（>=1 时再遇到 ZIP 直接跳过）。
ZIP_MAX_DEPTH = 1

# 直接按 utf-8 读的纯文本扩展名白名单。
PLAIN_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".log",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".csv", ".tsv",
    ".py", ".js", ".jsx", ".ts", ".tsx",
    ".java", ".c", ".cc", ".cpp", ".h", ".hpp",
    ".go", ".rs", ".rb", ".php", ".sh", ".bash", ".zsh",
    ".html", ".htm", ".xml", ".sql",
}


def _truncate(text: str, limit: int = PER_FILE_CHAR_LIMIT) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n…[内容过长，已截断，仅保留前 {limit} 字符]"


def _ext(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower()


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "[抽取失败: 未安装 pypdf，请在后端 requirements 安装后重试]"
    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as exc:
        return f"[抽取失败: PDF 解析错误: {exc}]"
    chunks = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            txt = page.extract_text() or ""
        except Exception as exc:
            txt = f"[第 {idx} 页解析失败: {exc}]"
        if txt.strip():
            chunks.append(f"## 第 {idx} 页\n{txt.strip()}")
    if not chunks:
        return "[PDF 中未抽取到可读文本，可能是扫描版或加密文件]"
    return "\n\n".join(chunks)


def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        return "[抽取失败: 未安装 python-docx]"
    try:
        doc = Document(io.BytesIO(data))
    except Exception as exc:
        return f"[抽取失败: docx 解析错误: {exc}]"
    parts = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            cells = [c for c in cells if c]
            if cells:
                parts.append(" | ".join(cells))
    if not parts:
        return "[docx 中未抽取到文本]"
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    try:
        import openpyxl
    except ImportError:
        return "[抽取失败: 未安装 openpyxl]"
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    except Exception as exc:
        return f"[抽取失败: xlsx 解析错误: {exc}]"
    chunks = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            chunks.append(f"## Sheet: {sheet.title}\n" + "\n".join(rows))
    if not chunks:
        return "[xlsx 中未抽取到内容]"
    return "\n\n".join(chunks)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "[抽取失败: 未安装 python-pptx]"
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        return f"[抽取失败: pptx 解析错误: {exc}]"
    chunks = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if text:
                lines.append(text)
        if lines:
            chunks.append(f"## 第 {idx} 页\n" + "\n".join(lines))
    if not chunks:
        return "[pptx 中未抽取到文本]"
    return "\n\n".join(chunks)


def _extract_plain(data: bytes) -> str:
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return "[抽取失败: 无法识别文本编码]"


def _extract_zip(data: bytes, depth: int = 0) -> str:
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except Exception as exc:
        return f"[抽取失败: zip 解析错误: {exc}]"
    chunks = []
    skipped = []
    handled = 0
    for info in zf.infolist():
        if info.is_dir():
            continue
        if handled >= ZIP_MAX_INNER_FILES:
            skipped.append(info.filename)
            continue
        inner_name = info.filename
        inner_ext = _ext(inner_name)
        if inner_ext == ".zip" and depth >= ZIP_MAX_DEPTH:
            skipped.append(inner_name + "（嵌套 zip 已忽略）")
            continue
        try:
            inner_bytes = zf.read(info)
        except Exception as exc:
            chunks.append(f"### {inner_name}\n[读取失败: {exc}]")
            handled += 1
            continue
        inner_text = _dispatch(inner_name, inner_bytes, depth=depth + 1)
        if inner_text is None:
            skipped.append(inner_name + "（不支持的类型）")
            continue
        chunks.append(f"### {inner_name}\n{inner_text}")
        handled += 1
    if not chunks and not skipped:
        return "[zip 内没有可处理的文件]"
    out = "\n\n".join(chunks)
    if skipped:
        out += "\n\n[以下文件被忽略：" + ", ".join(skipped) + "]"
    return out


def _dispatch(filename: str, data: bytes, depth: int = 0):
    """根据扩展名分发；不支持的类型返回 None。"""
    ext = _ext(filename)
    if ext == ".pdf":
        return _extract_pdf(data)
    if ext == ".docx":
        return _extract_docx(data)
    if ext == ".doc":
        return "[抽取失败: 旧版 .doc 格式不支持，请另存为 .docx 后重试]"
    if ext == ".xlsx":
        return _extract_xlsx(data)
    if ext == ".pptx":
        return _extract_pptx(data)
    if ext == ".zip":
        return _extract_zip(data, depth=depth)
    if ext in PLAIN_TEXT_EXTS:
        return _extract_plain(data)
    return None


def extract_text(filename: str, data: bytes) -> Tuple[str, bool]:
    """抽取文件文本。

    Returns:
        (text, supported): supported=False 表示扩展名不在支持列表中，
        text 此时是给用户看的提示文本。
    """
    text = _dispatch(filename, data, depth=0)
    if text is None:
        return (f"[不支持的文件类型: {filename}]", False)
    return (_truncate(text), True)
